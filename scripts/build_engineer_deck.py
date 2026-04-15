"""
Build an engineer-facing deep-dive deck for the Fidemo project.

Audience: engineers who'll read, run, modify, or fix this pipeline. Includes
the gotchas + materialization internals + version-pin reasoning that the
stakeholder deck deliberately omits. Real SQL and command syntax shown.

Output: docs/fidemo_engineer_deepdive.pptx

Run:  venv_dbt_duckdb/bin/python scripts/build_engineer_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme (matches the stakeholder deck for a consistent visual identity, but
# leans heavier on dark backgrounds — engineers spend more time in IDEs.)
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x0E, 0x2A, 0x47)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
SAND = RGBColor(0xE9, 0xC4, 0x6A)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF1, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1B, 0x2B, 0x3A)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Menlo"  # Falls back to Consolas / Courier on Windows

OUT = Path(__file__).resolve().parent.parent / "docs" / "fidemo_engineer_deepdive.pptx"


# ---------------------------------------------------------------------------
# Helpers (duplicated from build_stakeholder_deck.py — intentional; lets the
# two decks evolve independently without coupling)
# ---------------------------------------------------------------------------

def _set_run(run, *, text, size=18, bold=False, color=NAVY, font=BODY_FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_paragraph(p, text, *, size=18, bold=False, color=NAVY,
                   align=PP_ALIGN.LEFT, font=BODY_FONT, level=0):
    p.alignment = align
    p.level = level
    if p.runs:
        for r in p.runs:
            r.text = ""
        run = p.runs[0]
    else:
        run = p.add_run()
    _set_run(run, text=text, size=size, bold=bold, color=color, font=font)


def add_title_bar(slide, *, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.18)
    _set_paragraph(tf.paragraphs[0], title,
                   size=28, bold=True, color=WHITE, font=TITLE_FONT)
    if subtitle:
        p = tf.add_paragraph()
        _set_paragraph(p, subtitle, size=14, bold=False,
                       color=LIGHT, font=BODY_FONT)


def add_footer(slide, page_num, total):
    fl = slide.shapes.add_textbox(
        Inches(0.4), Inches(7.0), Inches(8), Inches(0.3)
    )
    _set_paragraph(fl.text_frame.paragraphs[0],
                   "Fidemo · engineering deep-dive",
                   size=10, color=GREY)
    fr = slide.shapes.add_textbox(
        Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3)
    )
    _set_paragraph(fr.text_frame.paragraphs[0],
                   f"{page_num} / {total}",
                   size=10, color=GREY, align=PP_ALIGN.RIGHT)


def add_textbox(slide, *, left, top, width, height, body, size=16,
                color=NAVY, bold=False, align=PP_ALIGN.LEFT, font=BODY_FONT):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    lines = body if isinstance(body, list) else [body]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_paragraph(p, line, size=size, color=color,
                       bold=bold, align=align, font=font)
    return tb


def add_bullet_list(slide, *, left, top, width, height, items, size=16):
    """items: list of (text, level)."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = ("•  " if lvl == 0 else "–  ") + text
        _set_paragraph(p, bullet, size=size, level=lvl)
    return tb


def add_table(slide, *, left, top, width, height, header, rows,
              header_fill=TEAL, alt_row_fill=LIGHT, body_size=11):
    cols = len(header)
    rcount = len(rows) + 1
    tbl = slide.shapes.add_table(
        rcount, cols, Inches(left), Inches(top),
        Inches(width), Inches(height)
    ).table
    for c, h in enumerate(header):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.clear()
        _set_paragraph(cell.text_frame.paragraphs[0], h,
                       size=13, bold=True, color=WHITE)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_row_fill
            cell.text_frame.clear()
            _set_paragraph(cell.text_frame.paragraphs[0], str(val),
                           size=body_size, color=NAVY)
    return tbl


def add_code_block(slide, *, left, top, width, height, code, size=11):
    """Render a fixed-width code block on a dark background."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = CODE_BG
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.05),
        Inches(width - 0.2), Inches(height - 0.1)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    lines = code if isinstance(code, list) else code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_paragraph(p, line, size=size, color=CODE_FG, font=CODE_FONT)
    return bg


# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

TOTAL = 14


# ---- 1. Title --------------------------------------------------------------
def slide_title(idx):
    s = prs.slides.add_slide(blank_layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_textbox(s, left=1.0, top=2.2, width=11.3, height=1.5,
                body="Fidemo", size=72, bold=True, color=WHITE)
    add_textbox(s, left=1.0, top=3.4, width=11.3, height=0.8,
                body="Engineering deep-dive — internals, gotchas, version pins",
                size=24, color=SAND)
    add_textbox(s, left=1.0, top=4.4, width=11.3, height=0.6,
                body="Companion to README.md / ARCHITECTURE.md / CLAUDE.md",
                size=16, color=LIGHT)
    add_textbox(s, left=1.0, top=5.5, width=11.3, height=0.4,
                body="Audience: anyone who'll read, run, modify, or fix this pipeline",
                size=14, color=LIGHT)


# ---- 2. Goals & non-goals --------------------------------------------------
def slide_goals(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Goals & non-goals",
                  subtitle="What we set out to demonstrate, and what we deliberately deferred")
    add_table(s, left=0.5, top=1.4, width=12.3, height=4.8,
              header=["Goal", "Status"],
              rows=[
                  ["DuckDB + dbt for object-storage transformation",
                   "✅ Working — staging view + 2 bronze materializations"],
                  ["Native DuckDB → SQL Server push (no pandas roundtrip)",
                   "✅ via `mssql` community extension + custom mat"],
                  ["SCD2 dimension in SQL Server, dbt-managed",
                   "✅ via dbt snapshot --target sqlserver"],
                  ["Two lake formats compared end-to-end",
                   "✅ Hive Parquet + DuckLake, identical silver outputs"],
                  ["Reproducible dev environment, host & container parity",
                   "✅ Dockerfile + .devcontainer/devcontainer.json"],
                  ["Data-quality auditing of dedup decisions",
                   "✅ store_failures + rejects model + reconciliation test"],
                  ["Production-grade secrets management",
                   "❌ Demo creds in plaintext — first task on any fork"],
                  ["CI/scheduled execution",
                   "❌ Container is ready; no Actions/Cron wired yet"],
              ])
    add_footer(s, idx, TOTAL)


# ---- 3. Architecture (same diagram as stakeholder, denser annotation) -----
def slide_architecture(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Architecture",
                  subtitle="Six dbt models + 2 snapshots + 1 sister rejects model")

    items = [
        ("RAW (MinIO):  s3://informat/seedcsv/year=YYYY/month=MM/day=DD/*.txt", 0),
        ("STAGING (DuckDB view):  stg_scb_bulkfil  ←  read_csv via httpfs, hive partitioning, latin-1", 0),
        ("BRONZE-PARQUET (MinIO):  bronze_scb_bulkfil_parquet  ←  external mat, partition_by year/month/day", 0),
        ("BRONZE-DUCKLAKE (MinIO + SQLite catalog):  bronze_scb_bulkfil_ducklake  ←  custom `ducklake` mat", 0),
        ("SILVER LANDING ×2 (SQL Server):  scb_bulkfil_landing_from_{parquet,ducklake}  ←  custom `mssql_native` mat, with QUALIFY dedup", 0),
        ("SILVER REJECTS (SQL Server):  scb_bulkfil_dedup_rejects  ←  same `mssql_native` mat, the QUALIFY-losers", 0),
        ("SCD2 ×2 (SQL Server):  snap_scb_bulkfil_scd2{,_from_ducklake}  ←  dbt snapshot, --target sqlserver, strategy=check", 0),
        ("", 0),
        ("Two custom materializations live in fidemo/macros/. Two `dbt run` targets (dev = DuckDB + extensions, sqlserver = native pyodbc).", 0),
    ]
    add_bullet_list(s, left=0.7, top=1.3, width=12.0, height=5.7,
                    items=items, size=14)
    add_footer(s, idx, TOTAL)


# ---- 4. Custom materializations (mssql_native + ducklake) -----------------
def slide_materializations(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Custom materializations",
                  subtitle="~50 lines of Jinja each — first-class in dbt-duckdb's plugin model")

    add_textbox(s, left=0.5, top=1.3, width=6.0, height=0.4,
                body="mssql_native — DuckDB → SQL Server",
                size=16, bold=True, color=NAVY)
    add_code_block(s, left=0.5, top=1.7, width=6.0, height=4.2,
                   code=[
                       "{% materialization mssql_native, adapter='duckdb' %}",
                       "  -- one-time per session:",
                       "  INSTALL mssql FROM community;",
                       "  LOAD mssql;",
                       "  ATTACH IF NOT EXISTS '<conn>' AS ms (TYPE mssql);",
                       "",
                       "  -- per model:",
                       "  CREATE OR REPLACE TABLE",
                       "    ms.finance.<model_name>",
                       "  AS SELECT * FROM ( {{ sql }} ) q;",
                       "",
                       "  return relation database='ms', schema='finance'",
                       "{% endmaterialization %}",
                   ], size=11)

    add_textbox(s, left=6.8, top=1.3, width=6.0, height=0.4,
                body="ducklake — DuckDB → DuckLake on S3",
                size=16, bold=True, color=NAVY)
    add_code_block(s, left=6.8, top=1.7, width=6.0, height=4.2,
                   code=[
                       "{% materialization ducklake, adapter='duckdb' %}",
                       "  INSTALL ducklake; INSTALL sqlite;",
                       "  LOAD ducklake;  LOAD sqlite;",
                       "  ATTACH IF NOT EXISTS",
                       "    'ducklake:sqlite:<catalog>.sqlite' AS lake",
                       "    (DATA_PATH 's3://...',",
                       "     AUTOMATIC_MIGRATION true);",
                       "  CREATE SCHEMA IF NOT EXISTS lake.bronze;",
                       "",
                       "  CREATE OR REPLACE TABLE",
                       "    lake.bronze.<model_name>",
                       "  AS SELECT * FROM ( {{ sql }} ) q;",
                       "{% endmaterialization %}",
                   ], size=11)

    add_textbox(s, left=0.5, top=6.1, width=12.3, height=0.6,
                body="Both opt-in via {{ config(materialized='mssql_native') }} or 'ducklake'. Strategies: replace | truncate | append.",
                size=12, color=GREY)
    add_footer(s, idx, TOTAL)


# ---- 5. Single DuckDB session + ATTACH lifecycle --------------------------
def slide_single_session(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Single-session constraint",
                  subtitle="Why staging + bronze + landing must all run in ONE `dbt run` invocation")

    items = [
        ("ATTACH state lives only within a DuckDB process.", 0),
        ("The ducklake materialization runs `ATTACH 'ducklake:sqlite:...' AS lake` inside the model build.", 0),
        ("Downstream `scb_bulkfil_landing_from_ducklake.sql` reads `lake.bronze.<name>` — that resolution requires `lake` to still be attached.", 0),
        ("If we split the build into two `dbt run` calls (e.g. `make build-bronze` then `make load-landings`), the second invocation starts a fresh session, sees no `lake`, and the SELECT fails.", 0),
        ("", 0),
        ("Mitigation in Makefile:", 0),
    ]
    add_bullet_list(s, left=0.7, top=1.3, width=12.0, height=2.6,
                    items=items, size=14)

    add_code_block(s, left=0.7, top=4.6, width=11.9, height=1.9,
                   code=[
                       "load-scb-bulkfil:  # ALL in one dbt run",
                       "    dbt run --target dev --select \\",
                       "        stg_scb_bulkfil \\",
                       "        bronze_scb_bulkfil_parquet bronze_scb_bulkfil_ducklake \\",
                       "        scb_bulkfil_landing_from_parquet scb_bulkfil_landing_from_ducklake \\",
                       "        scb_bulkfil_dedup_rejects",
                   ], size=11)
    add_footer(s, idx, TOTAL)


# ---- 6. The disable_transactions discovery --------------------------------
def slide_transactions(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="DuckLake writes silently dropped under default dbt-duckdb",
                  subtitle="Empirical discovery: dbt reports OK, ducklake_snapshot_changes records nothing")

    items = [
        ("Symptom: dbt run reported `OK created sql ducklake model …`. Re-querying `lake.bronze.<table>` showed no table. Probing the SQLite catalog directly: zero new entries in `ducklake_snapshot_changes`.", 0),
        ("Cause: dbt-duckdb wraps every model invocation in BEGIN…COMMIT against the local DuckDB. The COMMIT does not propagate to attached non-DuckDB catalogs (DuckLake-over-SQLite in our case).", 0),
        ("Without an explicit DuckLake-side commit, the writes vanish when the connection closes.", 0),
        ("Fix: opt out of dbt-duckdb's transaction wrapping. Our materializations are single-statement CTAS; we don't need atomicity guarantees from dbt.", 0),
    ]
    add_bullet_list(s, left=0.7, top=1.3, width=12.0, height=4.0,
                    items=items, size=14)

    add_code_block(s, left=0.7, top=5.7, width=11.9, height=1.4,
                   code=[
                       "# fidemo/profiles.yml",
                       "fidemo:",
                       "  outputs:",
                       "    dev:",
                       "      type: duckdb",
                       "      disable_transactions: true   # ← required for DuckLake writes",
                   ], size=11)
    add_footer(s, idx, TOTAL)


# ---- 7. CSV encoding gotcha ----------------------------------------------
def slide_encoding(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="DuckDB CSV encoding — `latin-1`, not ICU",
                  subtitle="Most ICU encoding aliases mangle ASCII into fullwidth Unicode")

    items = [
        ("Source files are Swedish (Latin-9 byte-content). DuckDB's CSV reader advertises 300+ ICU encoding names. Plausible-looking aliases include ISO8859_15, windows-1252, 8859_15.", 0),
        ("All of those silently apply Unicode compatibility normalisation. ASCII `F` (U+0046) becomes fullwidth `Ｆ` (U+FF26) — column NAMES get rewritten, not just values.", 0),
        ("Reliable encodings are limited to: utf-8, utf-16, latin-1.", 0),
        ("Latin-1 is functionally equivalent to ISO-8859-15 for our content (divergence is only on €, Š, œ, Ÿ — none appear in company names or addresses).", 0),
    ]
    add_bullet_list(s, left=0.7, top=1.3, width=12.0, height=3.5,
                    items=items, size=14)

    add_code_block(s, left=0.7, top=5.4, width=11.9, height=1.7,
                   code=[
                       "# fidemo/models/staging/_sources.yml",
                       "external_location: |",
                       "  read_csv('s3://informat/seedcsv/...',",
                       "           delim='\\t', header=true,",
                       "           encoding='latin-1',           ← only safe choice",
                       "           hive_partitioning=true, all_varchar=true)",
                   ], size=11)
    add_footer(s, idx, TOTAL)


# ---- 8. SCD2 dedup contract + MERGE failure -------------------------------
def slide_scd2_dedup(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="SCD2 source contract: one row per business key",
                  subtitle="Otherwise SQL Server's MERGE refuses with cardinality error 8672")

    items = [
        ("Bronze legitimately has duplicates: each cross-delivery file may carry the same peorgnr at a new effective_date.", 0),
        ("First snapshot run failed:  '[42000] The MERGE statement attempted to UPDATE or DELETE the same row more than once.'", 0),
        ("Root cause: snapshot's MERGE matched one target row to N source rows (one per delivery) on `unique_key=peorgnr`. SQL Server refuses that cardinality.", 0),
        ("Fix: dedup at silver, not bronze. Bronze keeps full history; silver landing emits one winning row per key (latest effective_date wins). DuckDB QUALIFY does this in one line.", 0),
    ]
    add_bullet_list(s, left=0.7, top=1.3, width=12.0, height=3.5,
                    items=items, size=14)

    add_code_block(s, left=0.7, top=5.3, width=11.9, height=1.8,
                   code=[
                       "-- fidemo/models/exports/scb_bulkfil_landing_from_parquet.sql",
                       "select * from {{ ref('bronze_scb_bulkfil_parquet') }}",
                       "qualify row_number() over (",
                       "    partition by peorgnr",
                       "    order by effective_date desc",
                       ") = 1   -- SQL Server has no QUALIFY; this runs in DuckDB",
                   ], size=11)
    add_footer(s, idx, TOTAL)


# ---- 9. Duplicate-row audit pattern ---------------------------------------
def slide_audit_pattern(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Duplicate-row audit pattern",
                  subtitle="Three artefacts: warn-test, sister rejects model, reconciliation singular test")

    add_table(s, left=0.5, top=1.3, width=12.3, height=4.5,
              header=["Artefact", "Where", "Purpose"],
              rows=[
                  ["unique test on bronze.peorgnr (severity: warn, store_failures: true)",
                   "Auto-table in main_dbt_test__audit",
                   "Lightweight signal: which keys are duplicated? Refreshed every run."],
                  ["scb_bulkfil_dedup_rejects (mssql_native model)",
                   "fidemo.finance.scb_bulkfil_dedup_rejects",
                   "Full-row context for losing duplicates, with _dedup_rank column"],
                  ["dedup_reconciliation.sql (singular dbt test, severity: error)",
                   "tests/dedup_reconciliation.sql",
                   "Asserts bronze_n = winners_n + rejects_n every run"],
                  ["audit_helper.compare_relations analyses",
                   "fidemo/analyses/compare_*",
                   "Diff two bronze paths' silver outputs row-for-row (must be identical)"],
              ])

    add_textbox(s, left=0.5, top=6.1, width=12.3, height=0.6,
                body="Two duplicate flavours can be distinguished post-hoc: cross-delivery (normal SCD2 input) vs intra-delivery (upstream corruption — analyst conversation).",
                size=12, color=GREY)
    add_footer(s, idx, TOTAL)


# ---- 10. Version pin matrix ----------------------------------------------
def slide_pins(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Version pin matrix",
                  subtitle="Each pin documented; re-evaluate by re-running the probe")
    add_table(s, left=0.5, top=1.4, width=12.3, height=5.0,
              header=["Pin", "Why"],
              rows=[
                  ["duckdb == 1.5.1",
                   "Newest DuckDB where the `mssql` community extension is published for osx_arm64 + linux_amd64 + linux_arm64 (1.5.2+ HTTP 404 across all three)"],
                  ["dbt-sqlserver == 1.9.0",
                   "Highest published; older 1.3.x imports the removed dbt.clients.agate_helper.empty_table and crashes under dbt-core 1.10+"],
                  ["dbt-duckdb == 1.10.1",
                   "Tracks dbt-core 1.11.x; supports DuckDB 1.5.x"],
                  ["Python >= 3.10 (we use 3.12 via uv)",
                   "Some transitive deps (e.g. black>=25.12) require Python 3.10+"],
                  ["DuckLake catalog format auto-migrates",
                   "ATTACH (..., AUTOMATIC_MIGRATION true) — irreversible upgrade per major"],
                  ["Flyway: noarch tarball (container) / OS-arch tarball (host)",
                   "10.10.0 doesn't publish a linux-arm64 archive; noarch + system JRE works on all containers"],
              ])
    add_footer(s, idx, TOTAL)


# ---- 11. dbt show asymmetry ----------------------------------------------
def slide_dbt_show(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="`dbt show` against mssql_native is asymmetric",
                  subtitle="Use `--target dev --select` for source-side, `--target sqlserver --inline + source()` for written-side")

    items = [
        ("`mssql_native` is a DuckDB-only materialization (INSTALL + ATTACH + CTAS run in DuckDB).", 0),
        ("`dbt show --target sqlserver --select <model>` asks the sqlserver adapter to re-execute the model's SELECT, which reads `ref('bronze_…')` — a Parquet location in MinIO. SQL Server can't reach Parquet → error.", 0),
        ("All v2 silver tables are declared as sources in `_sources.yml` so `source('finance_landing','…')` reaches them from either target.", 0),
    ]
    add_bullet_list(s, left=0.7, top=1.3, width=12.0, height=2.5,
                    items=items, size=14)

    add_code_block(s, left=0.7, top=4.0, width=11.9, height=3.0,
                   code=[
                       "# DuckDB-side preview — re-runs the model's SELECT in DuckDB",
                       "dbt show --target dev --select scb_bulkfil_dedup_rejects --limit 20",
                       "",
                       "# SQL-Server-side state — read what was actually written",
                       "dbt show --target sqlserver --inline \\",
                       "  \"select peorgnr, _dedup_rank, foretagsnamn",
                       "   from {{ source('finance_landing','scb_bulkfil_dedup_rejects') }}",
                       "   order by _dedup_rank desc\" --limit 20",
                   ], size=11)
    add_footer(s, idx, TOTAL)


# ---- 12. Tooling stack ----------------------------------------------------
def slide_tooling(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Tooling stack",
                  subtitle="Convention: `uv` everywhere; isolate per-tool when adapter pins conflict")

    add_table(s, left=0.5, top=1.3, width=12.3, height=4.7,
              header=["Tool", "Install", "Use case"],
              rows=[
                  ["dbt show",
                   "(built-in)",
                   "Quick terminal preview; ad-hoc SQL via --inline + source()"],
                  ["VS Code dbt Power User",
                   "(baked into devcontainer.json)",
                   "Lineage DAG, compiled SQL, model preview grid"],
                  ["harlequin (+ harlequin-mssql)",
                   "uv pip install --python /opt/venv/bin/python harlequin harlequin-mssql",
                   "Terminal TUI for both DuckDB and SQL Server"],
                  ["dbt_audit_helper",
                   "packages.yml (already wired)",
                   "Row-level diff between two refs; analyses/compare_*.sql"],
                  ["data-diff",
                   "uv pip install --python /opt/venv/bin/python 'data-diff[mssql]'",
                   "Cross-database row-level diff CLI"],
                  ["recce",
                   "uv tool install recce  (isolated, won't conflict with dbt pins)",
                   "PR-review web UI for model diffs"],
              ])
    add_footer(s, idx, TOTAL)


# ---- 13. Followups / open issues ------------------------------------------
def slide_followups(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Followups",
                  subtitle="Concrete TODOs for the next sprint")
    items = [
        ("Wire the pipeline into a scheduled CI job (GitHub Actions, reusing the dev container image).", 0),
        ("Move secrets from Makefile defaults to env-var injection or a real secrets manager (`.env` + `docker compose --env-file` is the minimum bar).", 0),
        ("Decide one bronze format for production. DuckLake if compliance/time-travel matters; Parquet if max interop wins.", 0),
        ("Migrate DuckLake catalog to PostgreSQL when the demo SQLite catalog can't keep up with concurrent writers.", 0),
        ("Promote the `mssql_native` materialization from project-local Jinja macro to a shareable dbt-duckdb plugin (Python `BasePlugin` subclass).", 0),
        ("Add SCD2 coverage for the second source domain (collateral / real-estate) using the same medallion pattern.", 0),
        ("Replace the on-run-end dbt_artifacts upload with something more curated (or wire it into BI).", 0),
    ]
    add_bullet_list(s, left=0.7, top=1.3, width=12.0, height=5.7,
                    items=items, size=15)
    add_footer(s, idx, TOTAL)


# ---- 14. Q&A --------------------------------------------------------------
def slide_thanks(idx):
    s = prs.slides.add_slide(blank_layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_textbox(s, left=1.0, top=2.5, width=11.3, height=1.5,
                body="Questions?", size=72, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, left=1.0, top=4.0, width=11.3, height=0.5,
                body="ARCHITECTURE.md — design rationale and trade-offs",
                size=18, color=SAND, align=PP_ALIGN.CENTER)
    add_textbox(s, left=1.0, top=4.5, width=11.3, height=0.5,
                body="CLAUDE.md — operator guide, full gotchas catalogue (13 items)",
                size=18, color=SAND, align=PP_ALIGN.CENTER)
    add_textbox(s, left=1.0, top=5.0, width=11.3, height=0.5,
                body="README.md — quickstart, makefile cheat-sheet, investigate tools",
                size=18, color=SAND, align=PP_ALIGN.CENTER)


# ---- Build all -----------------------------------------------------------
slides = [
    slide_title, slide_goals, slide_architecture, slide_materializations,
    slide_single_session, slide_transactions, slide_encoding, slide_scd2_dedup,
    slide_audit_pattern, slide_pins, slide_dbt_show, slide_tooling,
    slide_followups, slide_thanks,
]
assert len(slides) == TOTAL, f"slide count drift: {len(slides)} vs TOTAL={TOTAL}"

for i, fn in enumerate(slides, start=1):
    fn(i)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Wrote {OUT} ({TOTAL} slides)")
