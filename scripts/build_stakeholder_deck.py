"""
Build a stakeholder/management overview deck for the Fidemo project.

Audience: non-engineering stakeholders. Emphasis on outcomes, value, proof
points; minimal jargon, no code, no Makefile commands, no gotchas.

Output: docs/fidemo_stakeholder_overview.pptx

Run:  venv_dbt_duckdb/bin/python scripts/build_stakeholder_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x0E, 0x2A, 0x47)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
SAND = RGBColor(0xE9, 0xC4, 0x6A)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF1, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

OUT = (
    Path(__file__).resolve().parent.parent / "docs" / "fidemo_stakeholder_overview.pptx"
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _set_run(run, *, text, size=18, bold=False, color=NAVY, font=BODY_FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_paragraph(
    p,
    text,
    *,
    size=18,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
    font=BODY_FONT,
    level=0,
):
    p.alignment = align
    p.level = level
    if p.runs:
        for r in p.runs:
            r.text = ""
        run = p.runs[0]
    else:
        run = p.add_run()
    _set_run(run, text=text, size=size, bold=bold, color=color, font=font)


def add_title_bar(slide, *, title, subtitle=None):
    # Coloured bar across the top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.18)
    _set_paragraph(
        tf.paragraphs[0], title, size=28, bold=True, color=WHITE, font=TITLE_FONT
    )
    if subtitle:
        p = tf.add_paragraph()
        _set_paragraph(p, subtitle, size=14, bold=False, color=LIGHT, font=BODY_FONT)


def add_footer(slide, page_num, total):
    fl = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(8), Inches(0.3))
    _set_paragraph(
        fl.text_frame.paragraphs[0],
        "Fidemo · DuckDB → SQL Server lakehouse with SCD2",
        size=10,
        color=GREY,
    )
    fr = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3))
    _set_paragraph(
        fr.text_frame.paragraphs[0],
        f"{page_num} / {total}",
        size=10,
        color=GREY,
        align=PP_ALIGN.RIGHT,
    )


def add_textbox(
    slide,
    *,
    left,
    top,
    width,
    height,
    body,
    size=16,
    color=NAVY,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    lines = body if isinstance(body, list) else [body]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_paragraph(p, line, size=size, color=color, bold=bold, align=align)
    return tb


def add_bullet_list(slide, *, left, top, width, height, items, size=18):
    """items: list of (text, level)."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = ("•  " if lvl == 0 else "–  ") + text
        _set_paragraph(p, bullet, size=size, level=lvl)
    return tb


def add_table(
    slide,
    *,
    left,
    top,
    width,
    height,
    header,
    rows,
    header_fill=TEAL,
    alt_row_fill=LIGHT,
):
    cols = len(header)
    rcount = len(rows) + 1
    tbl = slide.shapes.add_table(
        rcount, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    # header
    for c, h in enumerate(header):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.clear()
        _set_paragraph(
            cell.text_frame.paragraphs[0], h, size=14, bold=True, color=WHITE
        )
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_row_fill
            cell.text_frame.clear()
            _set_paragraph(cell.text_frame.paragraphs[0], str(val), size=12, color=NAVY)
    return tbl


def add_pipeline_box(
    slide, *, left, top, width, height, label, fill=TEAL, text_color=WHITE
):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    lines = label if isinstance(label, list) else [label]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_paragraph(
            p,
            line,
            size=11 if i > 0 else 13,
            bold=(i == 0),
            color=text_color,
            align=PP_ALIGN.CENTER,
        )
    return shp


def add_arrow(slide, *, x1, y1, x2, y2, color=NAVY):
    line = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True


# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

TOTAL = 14


# ---- 1. Title --------------------------------------------------------------
def slide_title(idx):
    s = prs.slides.add_slide(blank_layout)
    bg = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_textbox(
        s,
        left=1.0,
        top=2.5,
        width=11.3,
        height=1.5,
        body="Fidemo",
        size=72,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        s,
        left=1.0,
        top=3.7,
        width=11.3,
        height=0.8,
        body="A DuckDB → SQL Server medallion lakehouse with SCD2 history",
        size=24,
        color=SAND,
    )
    add_textbox(
        s,
        left=1.0,
        top=5.5,
        width=11.3,
        height=0.4,
        body="Stakeholder overview · 2026",
        size=16,
        color=LIGHT,
    )


# ---- 2. Executive summary --------------------------------------------------
def slide_summary(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s, title="Executive summary", subtitle="What we built, in three sentences"
    )
    items = [
        (
            "A reproducible data pipeline that ingests SCB company-register deliveries, transforms them in DuckDB against object storage (MinIO), and lands a versioned dimension into SQL Server.",
            0,
        ),
        (
            'Two parallel "bronze" lake formats (hive Parquet and DuckLake) prove the architecture is portable; one production path is selected on rollout.',
            0,
        ),
        (
            'SCD2 change tracking captures every revision per company over time — we can answer "what did SCB say about company X on date Y?" without keeping every raw delivery file in the warehouse.',
            0,
        ),
        (
            "Runs identically on a developer laptop and in a container — no platform-specific setup steps.",
            0,
        ),
    ]
    add_bullet_list(s, left=0.7, top=1.4, width=12.0, height=5.5, items=items, size=20)
    add_footer(s, idx, TOTAL)


# ---- 3. Why this matters (business framing) -------------------------------
def slide_problem(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s,
        title="The problem we solved",
        subtitle="Why a lakehouse + SCD2, not just a warehouse load",
    )
    items = [
        (
            "Source data arrives as periodic full-extract files (deliveries every few months).",
            0,
        ),
        (
            "Each delivery overlaps the previous — companies are renamed, addresses change, organisations dissolve.",
            0,
        ),
        (
            'Stakeholders need both "what does the data say today" AND "what did the data say last quarter for company X".',
            0,
        ),
        (
            "A simple overwrite-on-load loses history; manual archival of raw files defers the problem instead of solving it.",
            0,
        ),
        ("", 0),
        (
            "Our answer: keep raw immutable in object storage, derive a versioned dimension table in SQL Server with explicit valid-from / valid-to columns. Standard Kimball SCD2 — battle-tested for 30 years.",
            0,
        ),
    ]
    add_bullet_list(s, left=0.7, top=1.4, width=12.0, height=5.5, items=items, size=18)
    add_footer(s, idx, TOTAL)


# ---- 4. Pipeline diagram (medallion) --------------------------------------
def slide_pipeline(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s,
        title="Pipeline at a glance",
        subtitle="Raw files → two bronze lake formats → SQL Server SCD2",
    )

    # Layer labels (left)
    add_textbox(
        s,
        left=0.3,
        top=1.5,
        width=1.6,
        height=0.4,
        body="RAW",
        size=14,
        bold=True,
        color=GREY,
    )
    add_textbox(
        s,
        left=0.3,
        top=2.7,
        width=1.6,
        height=0.4,
        body="STAGING",
        size=14,
        bold=True,
        color=GREY,
    )
    add_textbox(
        s,
        left=0.3,
        top=3.7,
        width=1.6,
        height=0.4,
        body="BRONZE",
        size=14,
        bold=True,
        color=GREY,
    )
    add_textbox(
        s,
        left=0.3,
        top=5.0,
        width=1.6,
        height=0.4,
        body="SILVER",
        size=14,
        bold=True,
        color=GREY,
    )
    add_textbox(
        s,
        left=0.3,
        top=6.3,
        width=1.6,
        height=0.4,
        body="HISTORY",
        size=14,
        bold=True,
        color=GREY,
    )

    # Boxes
    add_pipeline_box(
        s,
        left=2.2,
        top=1.4,
        width=9.0,
        height=0.9,
        label=[
            "Raw delivery files in object storage (MinIO)",
            "year=YYYY/month=MM/day=DD/scb_bulkfil_*.txt",
        ],
        fill=GREY,
    )

    add_pipeline_box(
        s,
        left=4.5,
        top=2.6,
        width=4.5,
        height=0.9,
        label=["Staging view (DuckDB)", "type-cast, normalise, partition-aware"],
        fill=NAVY,
    )

    add_pipeline_box(
        s,
        left=2.2,
        top=3.6,
        width=4.4,
        height=1.1,
        label=["Bronze · hive Parquet", "Open format, every tool reads it"],
        fill=TEAL,
    )
    add_pipeline_box(
        s,
        left=6.8,
        top=3.6,
        width=4.4,
        height=1.1,
        label=["Bronze · DuckLake", "ACID + time travel + schema evolution"],
        fill=TEAL,
    )

    add_pipeline_box(
        s,
        left=2.2,
        top=4.9,
        width=4.4,
        height=1.1,
        label=["Silver landing (SQL Server)", "Latest version per company"],
        fill=SAND,
        text_color=NAVY,
    )
    add_pipeline_box(
        s,
        left=6.8,
        top=4.9,
        width=4.4,
        height=1.1,
        label=["Silver landing (SQL Server)", "Latest version per company"],
        fill=SAND,
        text_color=NAVY,
    )

    add_pipeline_box(
        s,
        left=4.0,
        top=6.2,
        width=5.4,
        height=1.0,
        label=["SCD2 history dimension (SQL Server)", "valid_from · valid_to · scd_id"],
        fill=CORAL,
    )

    # Arrows
    add_arrow(s, x1=6.7, y1=2.3, x2=6.7, y2=2.6)
    add_arrow(s, x1=4.4, y1=3.5, x2=4.4, y2=3.6)
    add_arrow(s, x1=9.0, y1=3.5, x2=9.0, y2=3.6)
    add_arrow(s, x1=4.4, y1=4.7, x2=4.4, y2=4.9)
    add_arrow(s, x1=9.0, y1=4.7, x2=9.0, y2=4.9)
    add_arrow(s, x1=4.4, y1=6.0, x2=5.5, y2=6.2)
    add_arrow(s, x1=9.0, y1=6.0, x2=7.9, y2=6.2)

    add_footer(s, idx, TOTAL)


# ---- 5. Two bronze variants (comparison) ----------------------------------
def slide_two_bronze(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s,
        title="Two bronze formats, same data",
        subtitle="So we can pick the right one before going to production",
    )
    add_table(
        s,
        left=0.5,
        top=1.4,
        width=12.3,
        height=4.5,
        header=["", "Hive Parquet", "DuckLake"],
        rows=[
            [
                "Storage",
                "Plain Parquet files in MinIO",
                "Parquet + ACID metadata catalog",
            ],
            ["Read by", "Every analytics tool", "DuckDB-native, growing ecosystem"],
            ["Time travel", "No", "Yes — query the lake as of a past date"],
            ["Schema evolution", "Manual", "Tracked in catalog, automatic"],
            [
                "Operational cost",
                "Lowest possible",
                "One small SQLite or PostgreSQL catalog DB",
            ],
            [
                "Best fit",
                "Maximum interoperability",
                "Compliance / audit / point-in-time analytics",
            ],
        ],
    )
    add_textbox(
        s,
        left=0.5,
        top=6.3,
        width=12.3,
        height=0.6,
        body="Both produce identical SQL Server output — verified row-for-row by the project's diff tooling.",
        size=14,
        color=GREY,
    )
    add_footer(s, idx, TOTAL)


# ---- 6. SCD2 explained ----------------------------------------------------
def slide_scd2(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s,
        title="SCD2 — Slowly Changing Dimension type 2",
        subtitle="Standard pattern for tracking history over time",
    )

    items = [
        ("Each row in the dimension represents a company at a point in time.", 0),
        (
            "Two extra columns: valid_from (when this version became current) and valid_to (when it was superseded; NULL = current).",
            0,
        ),
        (
            "On every refresh, dbt's snapshot mechanism diffs incoming rows against the existing history:",
            0,
        ),
        ("  unchanged rows → no action", 1),
        (
            "  changed rows → close out the old version (set valid_to), insert the new version (valid_to = NULL)",
            1,
        ),
        ("  new rows → insert as current", 1),
        (
            'Result: any analyst query can ask "what did we know about company X on 2026-02-15?" with a simple WHERE clause.',
            0,
        ),
        (
            "Built natively against SQL Server using dbt's snapshot mechanism — no custom MERGE, no procedural code.",
            0,
        ),
    ]
    add_bullet_list(s, left=0.7, top=1.4, width=12.0, height=5.6, items=items, size=17)
    add_footer(s, idx, TOTAL)


# ---- 7. Quality / dedup audit --------------------------------------------
def slide_audit(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s,
        title="Data quality — what happens to bad rows",
        subtitle="Garbage in is quarantined, not silently dropped",
    )

    items = [
        (
            "Duplicate detection on bronze: any company appearing more than once is logged to a quarantine table that any analyst can query.",
            0,
        ),
        ("Two distinct duplicate flavours are tracked separately:", 0),
        (
            "  Cross-delivery — same company in last quarter's file and this quarter's file. Normal SCD2 input; older version becomes a closed-out history row.",
            1,
        ),
        (
            "  Intra-delivery — same company twice in ONE file. Indicates upstream system corruption; an analyst conversation should follow.",
            1,
        ),
        (
            "Reconciliation test verifies on every run: bronze row count = silver row count + quarantine row count. No rows lost.",
            0,
        ),
        (
            "All quality observations are queryable from BI tools alongside the SCD2 dimension itself.",
            0,
        ),
    ]
    add_bullet_list(s, left=0.7, top=1.4, width=12.0, height=5.6, items=items, size=17)
    add_footer(s, idx, TOTAL)


# ---- 8. Reproducibility (host vs container) -------------------------------
def slide_reproducibility(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s,
        title="Reproducibility",
        subtitle="Same pipeline on a laptop, in a container, in CI",
    )
    add_table(
        s,
        left=0.5,
        top=1.4,
        width=12.3,
        height=4.0,
        header=["Run mode", "Setup steps", "Best for"],
        rows=[
            [
                "Developer laptop (host)",
                "One-time: install Python 3.12, ODBC driver, Docker. Then `make`-driven flow.",
                "Fast iteration, debugging",
            ],
            [
                "Container (Docker / VS Code Dev Containers)",
                "Zero host setup beyond Docker. All deps baked into a 1.5 GB image.",
                "Onboarding, demo, CI",
            ],
            [
                "Future: CI pipeline",
                "Reuses the same container image. No drift between local and CI.",
                "Scheduled refresh, regression checks",
            ],
        ],
    )
    add_textbox(
        s,
        left=0.5,
        top=5.7,
        width=12.3,
        height=1.0,
        body="The container path absorbs every platform-specific friction we hit during development (driver versions, encoding handling, package pins). New collaborators are productive in the time it takes to pull an image.",
        size=14,
        color=GREY,
    )
    add_footer(s, idx, TOTAL)


# ---- 9. Observability / investigation -------------------------------------
def slide_observability(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s, title="Observability", subtitle="What can we see, measure, and compare"
    )
    items = [
        (
            "Per-run dbt artifacts capture: which models built, how long they took, which tests passed/warned/failed, lineage of every transformation.",
            0,
        ),
        (
            "Built-in tooling for ad-hoc inspection: terminal preview of any model, a TUI browser for both DuckDB and SQL Server, and a row-by-row diff CLI.",
            0,
        ),
        (
            "Audit-helper-driven diff between the two bronze paths runs as one command — proves the architecture stays consistent across formats.",
            0,
        ),
        (
            "All model + test results published to a dbt_artifacts table on every run; downstream BI dashboards or alerts can subscribe.",
            0,
        ),
        (
            "VS Code dev container ships with the dbt Power User extension preconfigured: lineage DAG, compiled SQL preview, and result browser one click away.",
            0,
        ),
    ]
    add_bullet_list(s, left=0.7, top=1.4, width=12.0, height=5.6, items=items, size=17)
    add_footer(s, idx, TOTAL)


# ---- 10. Security ---------------------------------------------------------
def slide_security(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s, title="Security posture", subtitle="Defense-in-depth, demo-grade by default"
    )
    items = [
        (
            "Least-privilege loader: the pipeline writes to SQL Server as a restricted user (`fidemo_loader`), scoped to the `finance` schema only. The system administrator account is used solely for one-time database creation.",
            0,
        ),
        (
            "Schema migrations are version-controlled and applied via Flyway — the database structure is reproducible and auditable, not configured by hand.",
            0,
        ),
        (
            "Corporate / self-signed certificates are bundled into a single CA bundle at install time and made available to every Python network call (uv, dbt, requests).",
            0,
        ),
        (
            "Demo credentials are documented in plain text inside the repository — appropriate for a sandbox; first step on any production fork is environment-variable injection or a real secrets manager.",
            0,
        ),
        (
            "Container image is built from official Microsoft + Astral upstream sources only — no third-party binaries.",
            0,
        ),
    ]
    add_bullet_list(s, left=0.7, top=1.4, width=12.0, height=5.6, items=items, size=16)
    add_footer(s, idx, TOTAL)


# ---- 11. Proof points / scorecard -----------------------------------------
def slide_scorecard(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(s, title="Proof points", subtitle="What's in the pipeline today")

    add_table(
        s,
        left=0.5,
        top=1.4,
        width=12.3,
        height=4.5,
        header=["Metric", "Status"],
        rows=[
            [
                "End-to-end pipeline runs",
                "✅ Passing — 6 of 6 dbt models, 2 of 2 snapshots",
            ],
            ["Records under SCD2 management", "198 companies, ready to grow"],
            [
                "Bronze formats verified equivalent",
                "✅ Both paths produce bit-identical SQL Server output",
            ],
            [
                "Quality tests",
                "Bronze uniqueness + null checks + reconciliation — all green",
            ],
            [
                "Documentation",
                "README + ARCHITECTURE + CLAUDE (operator guide), all current",
            ],
            [
                "Reproducibility",
                "Container build verified on macOS arm64; Linux x64/arm64 supported",
            ],
            [
                "Time to a fresh end-to-end run from `git clone`",
                "~5 minutes (mostly image build); pipeline itself executes in seconds",
            ],
        ],
    )
    add_footer(s, idx, TOTAL)


# ---- 12. Roadmap ----------------------------------------------------------
def slide_roadmap(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s, title="Open items / next steps", subtitle="What we'd do with the next sprint"
    )
    items = [
        (
            "Pick a single bronze format for production and retire the other side of the comparison.",
            0,
        ),
        (
            "Move DuckLake catalog from SQLite to PostgreSQL if multi-writer concurrency becomes a need.",
            0,
        ),
        (
            "Wire the pipeline into a scheduled CI job — same container image, automatic runs on new deliveries.",
            0,
        ),
        (
            "Replace plain-text credentials with a secrets manager (Vault, AWS Secrets Manager, or environment-variable injection).",
            0,
        ),
        (
            "Expand SCD2 coverage to additional source tables (collateral, real-estate, organisational hierarchy) using the same pattern.",
            0,
        ),
        (
            "Build BI dashboards on top of the SCD2 dimension — the heavy lifting is done; reporting is now configuration.",
            0,
        ),
    ]
    add_bullet_list(s, left=0.7, top=1.4, width=12.0, height=5.6, items=items, size=18)
    add_footer(s, idx, TOTAL)


# ---- 13. Risks ------------------------------------------------------------
def slide_risks(idx):
    s = prs.slides.add_slide(blank_layout)
    add_title_bar(
        s, title="Risks worth flagging", subtitle="What to watch as the pipeline scales"
    )
    add_table(
        s,
        left=0.5,
        top=1.4,
        width=12.3,
        height=4.5,
        header=["Risk", "Mitigation"],
        rows=[
            [
                "Single-host DuckLake catalog (SQLite) is single-writer.",
                "Swap to PostgreSQL catalog when concurrent writes are needed.",
            ],
            [
                "Community DuckDB extension for SQL Server is newer than the rest of the stack.",
                "Pinned to a verified working version. Re-evaluate on each upgrade.",
            ],
            [
                "Source files arrive in a non-Unicode encoding; misreading them silently corrupts data.",
                "Encoding is locked in source config and visually verified during ingest.",
            ],
            [
                "Demo credentials in source control if the repo is ever made public.",
                "First step on any non-demo fork: secrets manager. Documented in README.",
            ],
        ],
    )
    add_footer(s, idx, TOTAL)


# ---- 14. Q&A --------------------------------------------------------------
def slide_thanks(idx):
    s = prs.slides.add_slide(blank_layout)
    bg = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_textbox(
        s,
        left=1.0,
        top=2.7,
        width=11.3,
        height=1.5,
        body="Questions?",
        size=72,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        left=1.0,
        top=4.2,
        width=11.3,
        height=0.8,
        body="Detailed walk-through in README.md and ARCHITECTURE.md",
        size=20,
        color=SAND,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        left=1.0,
        top=5.0,
        width=11.3,
        height=0.4,
        body="Operator guide (commands, gotchas) in CLAUDE.md",
        size=16,
        color=LIGHT,
        align=PP_ALIGN.CENTER,
    )


# ---- Build all -----------------------------------------------------------
slides = [
    slide_title,
    slide_summary,
    slide_problem,
    slide_pipeline,
    slide_two_bronze,
    slide_scd2,
    slide_audit,
    slide_reproducibility,
    slide_observability,
    slide_security,
    slide_scorecard,
    slide_roadmap,
    slide_risks,
    slide_thanks,
]
assert len(slides) == TOTAL, f"slide count drift: {len(slides)} vs TOTAL={TOTAL}"

for i, fn in enumerate(slides, start=1):
    fn(i)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Wrote {OUT} ({TOTAL} slides)")
